import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Custom slide size (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Colors
    bg_color = RGBColor(10, 10, 10)
    coke_red = RGBColor(244, 0, 9)
    white = RGBColor(255, 255, 255)

    # Helper function to create standard layout
    def add_slide(title_text, content_bullets, visual_text=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_color

        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(45)
        p.font.bold = True
        p.font.color.rgb = coke_red

        # Content bullets
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(7.5), Inches(5))
        tf_content = content_box.text_frame
        tf_content.word_wrap = True
        for i, bullet in enumerate(content_bullets):
            p = tf_content.add_paragraph() if i > 0 else tf_content.paragraphs[0]
            p.text = bullet
            p.font.size = Pt(24)
            p.font.color.rgb = white
            p.space_after = Pt(14)
            p.level = 0
            
        # Optional Visual Box
        if visual_text:
            visual_box = slide.shapes.add_textbox(Inches(9), Inches(2), Inches(3.5), Inches(4))
            visual_box.fill.solid()
            visual_box.fill.fore_color.rgb = RGBColor(30, 0, 0) # Dark red box
            tf_vis = visual_box.text_frame
            tf_vis.word_wrap = True
            p = tf_vis.paragraphs[0]
            p.text = visual_text
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = coke_red
            p.alignment = PP_ALIGN.CENTER
            
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "COCA-COLA MANTRA 2026"
    p.font.size = Pt(70)
    p.font.bold = True
    p.font.color.rgb = coke_red
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Iconic Marketing: The Next Century"
    p2.font.size = Pt(36)
    p2.font.color.rgb = white
    p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph()
    p3.text = "\n#CokeAmbassador"
    p3.font.size = Pt(25)
    p3.font.color.rgb = coke_red
    p3.alignment = PP_ALIGN.CENTER

    # Slide 2: Hilltop
    add_slide(
        "1971: The Hilltop ➔ The Metaverse Concert",
        [
            "Strategic Vision: Host a synchronized, global virtual concert across Metaverse platforms like Fortnite and Roblox.",
            "Cultural Relevance: Gen Z and Alpha socialize in digital spaces. A digital Coke becomes the universal emote for connection.",
            "Execution: Users unlock exclusive concert access and digital wearables by scanning physical Coca-Cola cans.",
            "Brand Message: Real Magic transcends physical boundaries, uniting diverse cultures in a shared digital harmony."
        ],
        "Target Engagement:\n50M+ Concurrent Virtual Users\n\nPhygital Conversion:\n22% Scan-to-digital redemption"
    )

    # Slide 3: Share a Coke
    add_slide(
        "2011: Share a Coke ➔ AI Digital Twins",
        [
            "Strategic Vision: Shift from printed names to dynamic, generative AI personas linked to user profiles.",
            "Cultural Relevance: Driven by the rise of personalized AI and identity-driven tech.",
            "Execution: NFC-enabled cans open a web app where a holographic 'Coke Companion' shares playlists and vibes tailored to the user's mood.",
            "Brand Message: Coca-Cola doesn't just know your name; it understands your vibe."
        ],
        "Data Integration:\nAPI-First Spotify & Social Sync\n\nRetention Rate:\n3x Higher Engagement"
    )

    # Slide 4: Polar Bears
    add_slide(
        "1993: Polar Bears ➔ Arctic Guardians",
        [
            "Strategic Vision: Gamify recycling and environmental impact tracking through a mobile AR experience.",
            "Cultural Relevance: Modern consumers favor purpose-driven brands tackling climate change.",
            "Execution: Scanning a 100% rPET bottle drops an AR Polar Bear into the user's environment. Users 'feed' the bear by logging recycled bottles.",
            "Real-world Impact: Unlocking in-app achievements funds real-world donations to global ecological funds."
        ],
        "ESG Target:\n100% Collection & Recycling\n\nUser Action:\nPlay 2 Save"
    )

    # Slide 5: The Ambassador
    add_slide(
        "The 2026 Ambassador Vision",
        [
            "The Mantra: Respect the heritage, disrupt the delivery.",
            "The Goal: To weave Coca-Cola seamlessly into the cultural fabric of tomorrow—through Web3, AI, and uncompromising sustainability.",
            "The Promise: Driving measurable impact by turning passive consumers into active, engaged communities."
        ],
        "\n\nLet's Create the Future, Together.\n\nTaste the Future."
    )

    prs.save('Neon_Nostalgia_Presentation.pptx')
    print("Presentation created successfully as 'Neon_Nostalgia_Presentation.pptx'")

if __name__ == '__main__':
    create_presentation()
